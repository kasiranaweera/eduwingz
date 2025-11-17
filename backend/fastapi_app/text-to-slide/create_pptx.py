# 3_create_pptx.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import textwrap

def create_slideshow(title: str, bullets_list: list[list[str]], image_paths: list[str], output_name: str = "Final_Slides.pptx"):
    prs = Presentation()  # or Presentation("template.pptx") if you have one
    
    for i, (bullets, img_path) in enumerate(zip(bullets_list, image_paths)):
        slide_layout = prs.slide_layouts[5]  # Blank layout for full control
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf = title_shape.text_frame
        tf.text = f"{title} – Slide {i+1}"
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        
        # Image (left side)
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=Inches(5.5))
        
        # Bullets (right side)
        textbox = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(6), Inches(5.5))
        tf = textbox.text_frame
        tf.word_wrap = True
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.space_after = Pt(8)
            p.level = 0
    
    save_path = Path("outputs") / output_name
    prs.save(save_path)
    return str(save_path)

if __name__ == "__main__":
    bullets = [["• Sunlight energy", "• CO2 + H2O → Glucose + O2"], ["• Chlorophyll role", "• Light & dark reactions"]]
    images = ["outputs/images/slide_01.png", "outputs/images/slide_02.png"]
    create_slideshow("Photosynthesis", bullets, images)